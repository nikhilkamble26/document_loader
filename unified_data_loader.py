import os
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract


class UnifiedDataLoader:
    """
    Loads and extracts text from multiple file formats.
    Returns a list of documents with text and metadata.
    """

    def load(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            return self._load_pdf(file_path)
        elif ext in [".doc", ".docx"]:
            return self._load_docx(file_path)
        elif ext in [".ppt", ".pptx"]:
            return self._load_ppt(file_path)
        elif ext in [".xls", ".xlsx"]:
            return self._load_excel(file_path)
        elif ext in [".jpg", ".jpeg", ".png"]:
            return self._load_image(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def _load_pdf(self, file_path):
        reader = PdfReader(file_path)
        documents = []

        for page_num, page in enumerate(reader.pages, start=1):
            documents.append({
                "text": page.extract_text() or "",
                "metadata": {
                    "source": file_path,
                    "type": "pdf",
                    "page": page_num
                }
            })
        return documents

    def _load_docx(self, file_path):
        doc = Document(file_path)
        text = "\n".join(p.text for p in doc.paragraphs)

        return [{
            "text": text,
            "metadata": {
                "source": file_path,
                "type": "docx"
            }
        }]

    def _load_ppt(self, file_path):
        presentation = Presentation(file_path)
        documents = []

        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            documents.append({
                "text": "\n".join(slide_text),
                "metadata": {
                    "source": file_path,
                    "type": "pptx",
                    "slide": slide_num
                }
            })
        return documents

    def _load_excel(self, file_path):
        sheets = pd.read_excel(file_path, sheet_name=None)
        documents = []

        for sheet_name, df in sheets.items():
            documents.append({
                "text": df.to_string(index=False),
                "metadata": {
                    "source": file_path,
                    "type": "excel",
                    "sheet": sheet_name
                }
            })
        return documents

    def _load_image(self, file_path):
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)

        return [{
            "text": text,
            "metadata": {
                "source": file_path,
                "type": "image"
            }
        }]


if __name__ == "__main__":
    loader = UnifiedDataLoader()
    docs = loader.load("hdfc_doc.pdf")
    for doc in docs:
        print(doc["metadata"])
        print(doc["text"][:500])  # preview
